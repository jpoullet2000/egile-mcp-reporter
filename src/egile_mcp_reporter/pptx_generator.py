"""PowerPoint presentation generation using python-pptx."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


class PPTXGenerator:
    """Generate PowerPoint presentations."""

    def __init__(self):
        """Initialize PowerPoint generator."""
        pass

    def create_report(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        output_path: Path,
        **kwargs,
    ) -> Dict[str, Any]:
        """
        Create a PowerPoint presentation from structured sections.

        Args:
            title: Presentation title
            sections: List of section dictionaries
            output_path: Output PPTX path
            **kwargs: Additional options

        Returns:
            Dictionary with file metadata
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        import datetime

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1]
        subtitle.text = f"Generated: {datetime.datetime.now().strftime('%B %d, %Y')}"

        # Add content slides
        for section in sections:
            self._add_section_slide(prs, section)

        prs.save(str(output_path))

        return {
            "status": "success",
            "format": "pptx",
            "output_path": str(output_path),
            "size": output_path.stat().st_size,
            "title": title,
            "slides": len(prs.slides),
        }

    def from_markdown(
        self, markdown: str, output_path: str, title: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Convert markdown to PowerPoint (## = new slide).

        Args:
            markdown: Markdown content
            output_path: Output PPTX path
            title: Presentation title

        Returns:
            Dictionary with file metadata
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import datetime

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Parse markdown into slides
        slides_content = []
        current_slide = {"title": title or "Presentation", "content": []}

        for line in markdown.split("\n"):
            if line.startswith("# "):
                # Main title
                if current_slide["content"] or current_slide["title"] != title:
                    slides_content.append(current_slide)
                current_slide = {"title": line[2:].strip(), "content": []}
            elif line.startswith("## "):
                # New slide
                if current_slide["content"]:
                    slides_content.append(current_slide)
                current_slide = {"title": line[3:].strip(), "content": []}
            elif line.strip():
                current_slide["content"].append(line)

        if current_slide["content"] or len(slides_content) == 0:
            slides_content.append(current_slide)

        # Create slides
        for idx, slide_data in enumerate(slides_content):
            if idx == 0:
                # Title slide
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data["title"]
                if slide_data["content"]:
                    slide.placeholders[1].text = "\n".join(slide_data["content"][:3])
            else:
                # Content slide
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data["title"]
                
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()
                
                for line in slide_data["content"]:
                    p = text_frame.add_paragraph()
                    p.text = line.lstrip("- *").strip()
                    p.level = 0 if not line.startswith(("  ", "\t", "-", "*")) else 1

        prs.save(output_path)

        return {
            "status": "success",
            "format": "pptx",
            "output_path": str(output_path),
            "size": Path(output_path).stat().st_size,
            "slides": len(prs.slides),
        }

    def _add_section_slide(self, prs, section: Dict[str, Any]) -> None:
        """Add a slide for a section."""
        from pptx.util import Inches, Pt

        section_type = section.get("type", "text")
        section_title = section.get("title", "")

        if section_type == "text":
            # Content slide with title and text
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section_title

            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            
            content = section.get("content", "")
            for line in content.split("\n"):
                if line.strip():
                    p = text_frame.add_paragraph()
                    p.text = line.strip()

        elif section_type == "table":
            self._add_table_slide(prs, section)

        elif section_type == "list":
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section_title

            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()

            for item in section.get("items", []):
                p = text_frame.add_paragraph()
                p.text = str(item)
                p.level = 0

    def _add_table_slide(self, prs, section: Dict[str, Any]) -> None:
        """Add a slide with a table."""
        from pptx.util import Inches

        data = section.get("data", [])
        if not data:
            return

        columns = section.get("columns") or list(data[0].keys())

        # Blank slide
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_shape = slide.shapes.title
        title_shape.text = section.get("title", "Table")

        # Add table
        rows = len(data) + 1  # +1 for header
        cols = len(columns)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header row
        for col_idx, col_name in enumerate(columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            # Style header
            cell.fill.solid()
            cell.fill.fore_color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(
                30, 64, 175
            )
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.bold = True
            paragraph.font.size = __import__("pptx.util").util.Pt(11)
            paragraph.font.color.rgb = __import__("pptx.dml.color").dml.color.RGBColor(
                255, 255, 255
            )

        # Data rows
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, col_name in enumerate(columns):
                cell = table.cell(row_idx, col_idx)
                value = row_data.get(col_name, "")
                if isinstance(value, float):
                    cell.text = f"{value:,.2f}"
                else:
                    cell.text = str(value)
